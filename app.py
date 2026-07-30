import os
import io
import gc
import json
import re
import time
import traceback
import socket
import threading
import uuid
import base64

socket.setdefaulttimeout(180)

import requests as http_requests
from flask import Flask, render_template, request, send_file, jsonify
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

load_dotenv()

app = Flask(__name__)

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# ─── Job store ────────────────────────────────────────────────────────────────
_jobs = {}
_jobs_lock = threading.Lock()

def _store_job(job_id, data):
    with _jobs_lock:
        _jobs[job_id] = data

def _get_job(job_id):
    with _jobs_lock:
        return _jobs.get(job_id)

def _pop_job(job_id):
    with _jobs_lock:
        return _jobs.pop(job_id, None)

def _cleanup_jobs():
    cutoff = time.time() - 600
    with _jobs_lock:
        stale = [k for k, v in _jobs.items() if v.get("ts", 0) < cutoff]
        for k in stale:
            del _jobs[k]

# ─── Palette ──────────────────────────────────────────────────────────────────

C = {
    "primary":   RGBColor(0x2B, 0x45, 0x90),
    "secondary": RGBColor(0x4A, 0x90, 0xD9),
    "accent":    RGBColor(0xF5, 0xA6, 0x23),
    "dark":      RGBColor(0x2C, 0x2C, 0x2C),
    "white":     RGBColor(0xFF, 0xFF, 0xFF),
    "light":     RGBColor(0xF5, 0xF7, 0xFA),
    "success":   RGBColor(0x27, 0xAE, 0x60),
    "purple":    RGBColor(0x6C, 0x3D, 0xC4),
    "panel":     RGBColor(0xE8, 0xEF, 0xFB),
    "darkblue":  RGBColor(0x1A, 0x2F, 0x6B),
    "midblue":   RGBColor(0x3B, 0x5A, 0xAA),
    "lightblue": RGBColor(0xB8, 0xCC, 0xF0),
}

W = Inches(13.33)
H = Inches(7.5)

# ─── AI Content Generation ────────────────────────────────────────────────────

GEMINI_REST_URL = "https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={key}"


_IMAGE_KEYWORDS = ("image", "picture", "photo", "visual", "illustration", "graphic", "figure")

def _wants_images(additional_info):
    if not additional_info:
        return False
    lower = additional_info.lower()
    return any(kw in lower for kw in _IMAGE_KEYWORDS)


def _try_generate_image(prompt):
    """Call Gemini image generation. Returns PNG/JPEG bytes, or None on failure."""
    try:
        url = GEMINI_REST_URL.format(model="gemini-3.1-flash-image", key=GEMINI_API_KEY)
        payload = {
            "contents": [{"parts": [{"text": (
                f"Create a clean, colorful, flat educational illustration: {prompt}. "
                "No text or labels inside the image. Simple, minimal style. "
                "Suitable as a visual aid on a presentation slide."
            )}]}],
            "generationConfig": {"responseModalities": ["IMAGE"]},
        }
        session = http_requests.Session()
        try:
            resp = session.post(url, json=payload, timeout=30, headers={"Connection": "close"})
        finally:
            session.close()
        if not resp.ok:
            app.logger.warning(f"[img] {resp.status_code}: {resp.text[:120]}")
            return None
        parts = resp.json().get("candidates", [{}])[0].get("content", {}).get("parts", [])
        for part in parts:
            if "inlineData" in part:
                return base64.b64decode(part["inlineData"]["data"])
        return None
    except Exception as e:
        app.logger.warning(f"[img] generation failed: {e}")
        return None


def _call_gemini_rest(prompt, model):
    url = GEMINI_REST_URL.format(model=model, key=GEMINI_API_KEY)
    payload = {"contents": [{"parts": [{"text": prompt}]}]}
    t0 = time.time()
    session = http_requests.Session()
    try:
        resp = session.post(url, json=payload, timeout=170, headers={"Connection": "close"})
    finally:
        session.close()
    elapsed = round(time.time() - t0, 1)
    app.logger.info(f"[gemini] model={model} status={resp.status_code} elapsed={elapsed}s")
    if resp.status_code == 429:
        raise Exception(f"429 RESOURCE_EXHAUSTED {resp.text}")
    resp.raise_for_status()
    return resp.json()["candidates"][0]["content"]["parts"][0]["text"]


def generate_ppt_content(topic, audience, num_slides, additional_info):
    add_image_field = _wants_images(additional_info)
    image_field_schema = (
        '\n      "image_prompt": "concise 6-10 word visual description for a flat illustration of this slide\'s main concept",'
        if add_image_field else ""
    )

    prompt = f"""Create a comprehensive, engaging PowerPoint presentation for:

Topic: {topic}
Audience: {audience}
Content slides requested: {num_slides}
Extra context: {additional_info or "None"}

Return ONLY a JSON object (no markdown, no extra text) with this exact structure:
{{
  "title": "Catchy Presentation Title",
  "subtitle": "Descriptive subtitle in 8-12 words",
  "agenda": ["Topic area 1", "Topic area 2", "Topic area 3", "Topic area 4", "Topic area 5"],
  "slides": [
    {{
      "title": "Slide Title",
      "type": "content",
      "emoji": "📌",{image_field_schema}
      "key_points": [
        "Concise point, max 15 words",
        "Another clear point",
        "Third supporting point",
        "Fourth point if needed"
      ],
      "example": "One concrete real-world example that makes this tangible",
      "speaker_notes": "1-2 sentence presenter guide"
    }}
  ],
  "key_takeaways": [
    "Most important thing to remember",
    "Second key learning",
    "Third key learning",
    "Fourth key learning"
  ],
  "discussion_questions": [
    "Thought-provoking question 1?",
    "Question 2 that encourages reflection?",
    "Question 3?"
  ]
}}

Slide type rules:
- First 2 slides: type = "content" (foundational concepts)
- Every 3rd slide: type = "example" (real-world case study)
- One slide: type = "activity" (interactive exercise — also populate the "activity" field with a group exercise)
- Remaining slides: type = "content"
- Total: exactly {num_slides} slides in the "slides" array

Make it educational, clear, and suitable for {audience}.
Emoji should match the slide topic visually."""

    models = ["gemini-3.6-flash", "gemini-flash-lite-latest"]
    last_err = None
    for model in models:
        for attempt in range(2):
            try:
                text = _call_gemini_rest(prompt, model)
                text = text.strip()
                text = re.sub(r"^```(?:json)?\s*", "", text)
                text = re.sub(r"\s*```$", "", text)
                return json.loads(text)
            except Exception as e:
                last_err = e
                err_str = str(e)
                if "404" in err_str or "Not Found" in err_str:
                    break  # model not available — skip to next
                if "503" in err_str or "currently experiencing" in err_str:
                    if attempt == 0:
                        time.sleep(10)
                        continue
                    break
                if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
                    break  # skip to next model immediately — don't sleep
                raise  # non-rate-limit errors bubble up immediately
    err_str = str(last_err)
    if "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
        raise Exception("API quota exhausted. Please try again later or contact the administrator.")
    raise last_err

# ─── Slide Helpers ────────────────────────────────────────────────────────────

def _rect(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape



def _textbox(slide, x, y, w, h, text, size, color, bold=False, italic=False,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = align
    return tb


def _header(slide, title, bg=None, fg=None, height=Inches(1.4), font_size=26):
    bg = bg or C["primary"]
    fg = fg or C["white"]
    _rect(slide, 0, 0, W, height, bg)
    _textbox(slide, Inches(0.45), Inches(0.2), Inches(12.4), height - Inches(0.2),
             title, font_size, fg, bold=True)


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

# ─── Efficient multi-paragraph helper ─────────────────────────────────────────

def _bullets_textbox(slide, x, y, w, h, items, size, color, prefix="  •  ",
                     space_after=Pt(10), space_before=Pt(6), line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix}{item}"
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = space_after
        p.space_before = space_before
        if line_spacing:
            p.line_spacing = line_spacing
    return tb

# ─── Individual Slide Builders ────────────────────────────────────────────────

def _title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 4 shapes total
    _rect(slide, 0, 0, W, H, C["primary"])
    _rect(slide, 0, Inches(5.4), W, Inches(2.1), C["darkblue"])
    _rect(slide, Inches(1.0), Inches(5.25), Inches(11.33), Inches(0.07), C["accent"])
    _textbox(slide, Inches(1.0), Inches(1.4), Inches(10.5), Inches(2.8),
             data["title"], 44, C["white"], bold=True, wrap=True)
    _textbox(slide, Inches(1.0), Inches(4.3), Inches(10.5), Inches(0.9),
             data["subtitle"], 20, C["lightblue"], wrap=True)
    _textbox(slide, Inches(1.0), Inches(6.5), Inches(5.0), Inches(0.6),
             "✦ Interactive Presentation", 13, C["accent"], bold=True)


def _agenda_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["white"])
    _header(slide, "📋  What We'll Cover Today")
    _rect(slide, Inches(0.5), Inches(1.6), Inches(0.1), Inches(5.6), C["accent"])
    items = [f"{i+1}.   {item}" for i, item in enumerate(data.get("agenda", [])[:7])]
    n = max(1, len(items))
    fsize = Pt(24) if n <= 4 else Pt(21) if n <= 6 else Pt(18)
    spc   = Pt(18) if n <= 4 else Pt(14) if n <= 6 else Pt(10)
    _bullets_textbox(slide, Inches(1.0), Inches(1.65), Inches(11.5), Inches(5.6),
                     items, fsize, C["dark"], prefix="", space_after=spc, space_before=spc)
    _notes(slide, "Walk through the agenda. Let students know what they'll learn and why it matters.")


def _content_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["white"])
    _header(slide, f"{slide_data.get('emoji','📌')}  {slide_data['title']}")

    points      = slide_data.get("key_points", [])[:5]
    example     = slide_data.get("example", "")
    image_bytes = slide_data.get("image_bytes")
    n = max(1, len(points))

    if image_bytes:
        # Two-column: text left (7.2" wide), image right (5.1" wide)
        text_w = Inches(7.2)
        img_x  = Inches(7.7)
        img_w  = Inches(5.1)

        if n <= 2:
            fsize, spc = Pt(26), Pt(22)
        elif n == 3:
            fsize, spc = Pt(23), Pt(18)
        elif n == 4:
            fsize, spc = Pt(21), Pt(14)
        else:
            fsize, spc = Pt(19), Pt(10)

        if example:
            _bullets_textbox(slide, Inches(0.5), Inches(1.55), text_w, Inches(4.0),
                             points, fsize, C["dark"], space_after=spc, space_before=spc)
            _rect(slide, 0, Inches(5.65), Inches(7.5), Inches(1.75), C["panel"])
            _textbox(slide, Inches(0.55), Inches(5.8), Inches(6.9), Inches(1.5),
                     f"💡  {example}", 15, C["primary"], bold=False, italic=True, wrap=True)
        else:
            _bullets_textbox(slide, Inches(0.5), Inches(1.55), text_w, Inches(5.7),
                             points, fsize, C["dark"], space_after=spc, space_before=spc)

        try:
            slide.shapes.add_picture(io.BytesIO(image_bytes), img_x, Inches(1.6), img_w, Inches(5.2))
        except Exception as ex:
            app.logger.warning(f"[img] embed failed: {ex}")

    else:
        # Original full-width layout
        if n <= 2:
            fsize, spc = Pt(28), Pt(26)
        elif n == 3:
            fsize, spc = Pt(25), Pt(22)
        elif n == 4:
            fsize, spc = Pt(23), Pt(18)
        else:
            fsize, spc = Pt(21), Pt(14)

        if example:
            _bullets_textbox(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(4.2),
                             points, fsize, C["dark"], space_after=spc, space_before=spc)
            _rect(slide, 0, Inches(5.85), W, Inches(1.55), C["panel"])
            _textbox(slide, Inches(0.65), Inches(6.0), Inches(12.0), Inches(1.3),
                     f"💡  Example:   {example}", 18, C["primary"], bold=False, italic=True, wrap=True)
        else:
            _bullets_textbox(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(5.7),
                             points, fsize, C["dark"], space_after=spc, space_before=spc)

    _rect(slide, 0, Inches(7.4), W, Inches(0.1), C["secondary"])
    _notes(slide, slide_data.get("speaker_notes", ""))


def _example_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["light"])
    _header(slide, f"💡  {slide_data['title']} — Real-World Example", bg=C["secondary"])
    example_text = slide_data.get("example", "")
    # Large example panel — top half of body
    _rect(slide, Inches(0.45), Inches(1.55), Inches(12.4), Inches(3.3), C["panel"],
          line_color=C["secondary"])
    _textbox(slide, Inches(0.75), Inches(1.7), Inches(11.8), Inches(3.1),
             example_text, 20, C["dark"], wrap=True)
    # Takeaway bullets — bottom half
    points = slide_data.get("key_points", [])[:4]
    if points:
        n = len(points)
        fsize = Pt(21) if n <= 2 else Pt(19)
        spc   = Pt(16) if n <= 2 else Pt(12)
        _bullets_textbox(slide, Inches(0.6), Inches(5.05), Inches(12.1), Inches(2.25),
                         points, fsize, C["dark"], prefix="✅  ",
                         space_after=spc, space_before=spc)
    _notes(slide, slide_data.get("speaker_notes", ""))


def _activity_slide(prs, slide_data):
    # 5 shapes total (was 7)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["primary"])
    _textbox(slide, Inches(0.5), Inches(0.3), Inches(4.0), Inches(0.6),
             "🎯  ACTIVITY", 16, C["accent"], bold=True)
    _textbox(slide, Inches(0.5), Inches(1.0), Inches(12.0), Inches(1.4),
             slide_data["title"], 34, C["white"], bold=True, wrap=True)
    activity_text = slide_data.get("activity") or slide_data.get("example", "Activity coming soon.")
    _rect(slide, Inches(0.5), Inches(2.55), Inches(12.33), Inches(3.5), C["darkblue"])
    _textbox(slide, Inches(0.85), Inches(2.75), Inches(11.63), Inches(3.1),
             activity_text, 19, C["white"], wrap=True)
    _textbox(slide, Inches(0.5), Inches(6.3), Inches(12.0), Inches(0.7),
             "⏱  Take 3 minutes — then share with the group", 14, C["accent"])
    _notes(slide, slide_data.get("speaker_notes", ""))


def _summary_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["white"])
    _header(slide, "🎯  Key Takeaways", bg=C["success"])
    takeaways = [f"{i+1}.   {tk}" for i, tk in enumerate(data.get("key_takeaways", [])[:5])]
    n = max(1, len(takeaways))
    fsize = Pt(25) if n <= 3 else Pt(22) if n == 4 else Pt(20)
    spc   = Pt(22) if n <= 3 else Pt(16) if n == 4 else Pt(12)
    questions = data.get("discussion_questions", [])
    bullet_h = Inches(4.7) if questions else Inches(5.7)
    _bullets_textbox(slide, Inches(0.6), Inches(1.6), Inches(12.1), bullet_h,
                     takeaways, fsize, C["dark"], prefix="", space_after=spc, space_before=spc)
    if questions:
        _rect(slide, 0, Inches(6.5), W, Inches(1.0), C["panel"])
        _textbox(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.8),
                 f"💬  Discussion:   {questions[0]}", 17, C["primary"], italic=True, wrap=True)
    _notes(slide, "Summarize the key points. Ask students to reflect on what surprised them most.")


def _qa_slide(prs, data):
    # 4 shapes total (was 4 + N)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["primary"])
    _textbox(slide, Inches(1.5), Inches(1.5), Inches(9.0), Inches(2.5),
             "Questions?", 62, C["white"], bold=True, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(1.5), Inches(4.0), Inches(9.0), Inches(0.9),
             "Let's discuss! 🙋", 28, C["accent"], align=PP_ALIGN.CENTER)
    questions = data.get("discussion_questions", [])
    if questions:
        _bullets_textbox(slide, Inches(1.0), Inches(5.0), Inches(11.33), Inches(2.2),
                         questions[:3], Pt(18), C["lightblue"],
                         prefix="  •  ", space_after=Pt(14), space_before=Pt(8))
    _notes(slide, "Open the floor for questions. Use discussion prompts above to spark conversation if needed.")

# ─── Assemble Presentation ────────────────────────────────────────────────────

def create_presentation(data):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    _title_slide(prs, data)
    _agenda_slide(prs, data)

    for s in data.get("slides", []):
        t = s.get("type", "content")
        if t == "example":
            _example_slide(prs, s)
        elif t == "activity":
            _activity_slide(prs, s)
        else:
            _content_slide(prs, s)

    _summary_slide(prs, data)
    _qa_slide(prs, data)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ─── Routes ───────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/ping")
def ping():
    t0 = time.time()
    try:
        result = _call_gemini_rest("Reply with the single word: pong", "gemini-3.6-flash")
        return jsonify({"status": "ok", "reply": result.strip(), "elapsed_s": round(time.time() - t0, 1)})
    except Exception as e:
        return jsonify({"status": "error", "error": str(e)[:200], "elapsed_s": round(time.time() - t0, 1)}), 500


@app.route("/generate", methods=["POST"])
def generate():
    if not GEMINI_API_KEY:
        return jsonify({"error": "GEMINI_API_KEY is not set."}), 500

    topic = (request.form.get("topic") or "").strip()[:200]
    if not topic:
        return jsonify({"error": "Please enter a topic."}), 400

    audience = request.form.get("audience", "college students").strip()
    num_slides = max(4, min(int(request.form.get("num_slides", 10)), 30))
    additional_info = (request.form.get("additional_info") or "").strip()[:300]

    job_id = uuid.uuid4().hex
    _store_job(job_id, {"status": "pending", "ts": time.time()})

    def _run():
        try:
            t0 = time.time()
            ppt_data = generate_ppt_content(topic, audience, num_slides, additional_info)
            app.logger.info(f"[job:{job_id[:8]}] AI done in {round(time.time()-t0,1)}s")

            # Generate images for up to 3 content slides when user requested them
            if _wants_images(additional_info) and ppt_data.get("slides"):
                img_count = 0
                for slide in ppt_data["slides"]:
                    if img_count >= 3:
                        break
                    if slide.get("type") == "content" and slide.get("image_prompt"):
                        app.logger.info(f"[job:{job_id[:8]}] generating image: {slide['image_prompt'][:50]}")
                        img_bytes = _try_generate_image(slide["image_prompt"])
                        if img_bytes:
                            slide["image_bytes"] = img_bytes
                            img_count += 1
                app.logger.info(f"[job:{job_id[:8]}] images={img_count}")

            safe = re.sub(r"[^a-zA-Z0-9 ]", "", ppt_data.get("title", topic))[:40]
            filename = safe.strip().replace(" ", "_") + ".pptx"
            buf = create_presentation(ppt_data)
            del ppt_data
            gc.collect()
            app.logger.info(f"[job:{job_id[:8]}] total={round(time.time()-t0,1)}s")
            _store_job(job_id, {"status": "done", "data": buf.getvalue(), "filename": filename, "ts": time.time()})
        except Exception as e:
            traceback.print_exc()
            msg = str(e)
            if "429" in msg or "RESOURCE_EXHAUSTED" in msg:
                err = "API rate limit reached. Please wait a moment and try again."
            elif "quota" in msg.lower():
                err = "Daily API quota exceeded. Please try again tomorrow."
            elif "timed out" in msg.lower() or "timeout" in msg.lower():
                err = "The AI took too long to respond. Please try again with fewer slides."
            elif "json" in msg.lower():
                err = "AI returned an unexpected response. Please try again."
            else:
                err = "Something went wrong. Please try again."
            _store_job(job_id, {"status": "error", "error": err, "ts": time.time()})

    threading.Thread(target=_run, daemon=True).start()
    return jsonify({"job_id": job_id})


@app.route("/status/<job_id>")
def job_status(job_id):
    _cleanup_jobs()
    job = _get_job(job_id)
    if not job:
        return jsonify({"status": "not_found"}), 404
    return jsonify({"status": job["status"], "error": job.get("error", "")})


@app.route("/download/<job_id>")
def job_download(job_id):
    job = _pop_job(job_id)
    if not job or job["status"] != "done":
        return jsonify({"error": "File not ready or already downloaded. Please generate again."}), 404
    buf = io.BytesIO(job["data"])
    return send_file(
        buf,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=job.get("filename", "presentation.pptx"),
    )


@app.errorhandler(500)
def internal_error(e):
    return jsonify({"error": f"Internal server error: {str(e)}"}), 500


@app.errorhandler(Exception)
def unhandled_exception(e):
    traceback.print_exc()
    return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(debug=False, port=port)
